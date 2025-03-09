import streamlit as st
from PyPDF2 import PdfReader
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.cluster import KMeans
from pptx import Presentation
from pptx.util import Pt
import spacy
from spacy.lang.en.stop_words import STOP_WORDS
from string import punctuation
from heapq import nlargest
import os


# Function to summarize text
def summarize_text(text, ratio):
    nlp = spacy.load('en_core_web_sm')
    doc = nlp(text)
    word_frequencies = {}
    for word in doc:
        if word.text.lower() not in STOP_WORDS and word.text.lower() not in punctuation:
            word_frequencies[word.text.lower()] = word_frequencies.get(word.text.lower(), 0) + 1

    max_frequency = max(word_frequencies.values())
    for word in word_frequencies.keys():
        word_frequencies[word] /= max_frequency

    sentence_tokens = list(doc.sents)
    sentence_scores = {}
    for sent in sentence_tokens:
        for word in sent:
            if word.text.lower() in word_frequencies:
                sentence_scores[sent] = sentence_scores.get(sent, 0) + word_frequencies[word.text.lower()]

    select_length = max(1, int(len(sentence_tokens) * ratio))
    summary_sentences = nlargest(select_length, sentence_scores, key=sentence_scores.get)
    summary = ' '.join([sent.text for sent in summary_sentences])

    return summary


# Function to divide text into topics
def divide_into_topics(text, n_topics=3):
    sentences = text.split('. ')
    vectorizer = TfidfVectorizer(stop_words='english')
    X = vectorizer.fit_transform(sentences)

    kmeans = KMeans(n_clusters=n_topics, random_state=42)
    kmeans.fit(X)
    labels = kmeans.labels_

    keywords = []
    for i in range(n_topics):
        cluster_center = kmeans.cluster_centers_[i]
        top_indices = cluster_center.argsort()[-3:][::-1]
        keywords.append(", ".join([vectorizer.get_feature_names_out()[idx] for idx in top_indices]))

    topics = {}
    for i, label in enumerate(labels):
        topics.setdefault(label, []).append(sentences[i])

    return topics, keywords


# Function to create PowerPoint presentation
def create_presentation_with_organic_template(topics, topic_keywords,title, subtitle, template_file="organic_template.pptx", output_file="new.pptx"):
    # Load the existing template presentation with the "Organic" theme
    presentation = Presentation(template_file)

    # Slide layout: title slide
    slide_layout = presentation.slide_layouts[0]  # Title slide layout
    slide = presentation.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = input("Enter the title of summary: ").strip()
    subtitle.text = input("Enter the sub-title of summary: ").strip()

    # Content Slides for Each Topic
    slide_layout = presentation.slide_layouts[1]  # Title and content layout
    max_sentences = 5
    end_sentence_marker = "."
    first_slide_flag = True  # Used to distinguish the first slide of each topic

    for i, (topic_id, sentences) in enumerate(topics.items()):
        if not sentences:
            continue  # Skip empty topics

        full_text = '. '.join(sentences)  # Combine all sentences for the topic
        words = full_text.split()
        chunks = []  # Stores 5-sentence chunks for the topic

        sentence_count = 0
        temp_chunk = []

        # Create chunks with up to 5 sentences
        for word in words:
            temp_chunk.append(word)
            if end_sentence_marker in word:
                sentence_count += 1
            if sentence_count == max_sentences:
                chunks.append(' '.join(temp_chunk))
                temp_chunk = []
                sentence_count = 0

        # Add remaining words in the last chunk
        if temp_chunk:
            chunks.append(' '.join(temp_chunk))

        # Create slides for each chunk
        for chunk_idx, chunk in enumerate(chunks):
            if not chunk.strip():  # Skip blank chunks
                continue

            slide = presentation.slides.add_slide(slide_layout)
            title = slide.shapes.title
            content = slide.placeholders[1]

            # Title for the first slide of the topic
            if first_slide_flag:
                title.text = f"Topic: {topic_keywords[i]}"
                first_slide_flag = False
            else:
                title.text = f"Topic: {topic_keywords[i]} (continued)"

            # Set slide content
            content.text = chunk

            # Adjust font size for the content
            for paragraph in content.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(18)  # Set font size to 18 points

        # Reset the flag for the next topic
        first_slide_flag = True

    # Delete the starting slide
    xml_slides = presentation.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.remove(slides[0])
    presentation.save(output_file)
    # Save the presentation
    print(f"Presentation saved as {output_file}")

# Streamlit UI
def main():
    st.title("PDF to PowerPoint Converter")
    st.write("Upload a PDF file, and this tool will generate a PowerPoint presentation with summaries.")

    uploaded_file = st.file_uploader("Upload PDF", type=["pdf"])

    if uploaded_file is not None:
        st.write("File uploaded successfully.")
        summary_ratio = st.slider("Select summary ratio (e.g., 0.2 for 20%):", 0.1, 1.0, 0.2)
        n_topics = st.slider("Select number of topics to divide into:", 1, 10, 3)
        title = st.text_input("Enter the title for the presentation:")
        subtitle = st.text_input("Enter the subtitle for the presentation:")

        if st.button("Generate PowerPoint"):
            with st.spinner("Processing..."):
                try:
                    # Read PDF content
                    pdf_reader = PdfReader(uploaded_file)
                    full_text = ""
                    for page in pdf_reader.pages:
                        full_text += page.extract_text()

                    if not full_text.strip():
                        st.error("The PDF seems to be empty or unsupported.")
                        return

                    # Summarize and divide text
                    summary = summarize_text(full_text, summary_ratio)
                    topics, topic_keywords = divide_into_topics(summary, n_topics)

                    # Create PowerPoint presentation
                    output_file = "review_presentation.pptx"
                    create_presentation_with_organic_template(topics, topic_keywords, title, subtitle)

                    # Allow user to download the presentation
                    with open(output_file, "rb") as f:
                        st.download_button(
                            label="Download Presentation",
                            data=f,
                            file_name=output_file,
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        )
                    os.remove(output_file)  # Clean up file after download
                except Exception as e:
                    st.error(f"An error occurred: {e}")


if __name__ == "__main__":
    main()