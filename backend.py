# 5 sen per slide
import spacy
from spacy.lang.en.stop_words import STOP_WORDS 
from string import punctuation #to filter punctuations 
from heapq import nlargest #to select important sentences 
from PyPDF2 import PdfReader #read pdf
from sklearn.feature_extraction.text import TfidfVectorizer #text to numerical rep for clustering
from sklearn.cluster import KMeans #to create cluster of data
from pptx import Presentation #create presentation
from pptx.util import Pt  # Import for setting font size

# Function to summarize text
def summarize_text(text, ratio):
    nlp = spacy.load('en_core_web_sm')
    doc = nlp(text)
    
    word_frequencies = {}
    for word in doc:
        if word.text.lower() not in STOP_WORDS and word.text.lower() not in punctuation:
            word_frequencies[word.text.lower()] = word_frequencies.get(word.text.lower(), 0) + 1
    
    max_frequency = max(word_frequencies.values())
    for word in word_frequencies.keys():
        word_frequencies[word] /= max_frequency
    
    sentence_tokens = list(doc.sents)
    sentence_scores = {}
    for sent in sentence_tokens:
        for word in sent:
            if word.text.lower() in word_frequencies:
                sentence_scores[sent] = sentence_scores.get(sent, 0) + word_frequencies[word.text.lower()]
    
    select_length = max(1, int(len(sentence_tokens) * ratio))
    summary_sentences = nlargest(select_length, sentence_scores, key=sentence_scores.get)
    summary = ' '.join([sent.text for sent in summary_sentences])
    
    return summary

# Function to divide text into topics
def divide_into_topics(text, n_topics=3):
    sentences = text.split('. ')
    vectorizer = TfidfVectorizer(stop_words='english')
    X = vectorizer.fit_transform(sentences)
    
    kmeans = KMeans(n_clusters=n_topics, random_state=42)
    kmeans.fit(X)
    labels = kmeans.labels_
    
    # Extract topic keywords for slide titles
    keywords = []
    for i in range(n_topics):
        cluster_center = kmeans.cluster_centers_[i]
        top_indices = cluster_center.argsort()[-3:][::-1]
        keywords.append(", ".join([vectorizer.get_feature_names_out()[idx] for idx in top_indices]))
    
    # Group sentences into topics
    topics = {}
    for i, label in enumerate(labels):
        topics.setdefault(label, []).append(sentences[i])
    
    return topics, keywords

# Function to create PowerPoint from topics with adjusted font size
def create_presentation_with_topics(topics, topic_keywords, output_file="output_presentation.pptx"):
    presentation = Presentation()

    # Title Slide
    slide_layout = presentation.slide_layouts[0]  # Title slide layout
    slide = presentation.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = input("Enter the title of summary: ").strip()
    subtitle.text = input("Enter the sub-title of summary: ").strip()

    # Content Slides for Each Topic
    slide_layout = presentation.slide_layouts[1]  # Title and content layout
    max_sentences = 5
    end_sentence_marker = "."
    first_slide_flag = True  # Used to distinguish first slide of each topic

    for i, (topic_id, sentences) in enumerate(topics.items()):
        if not sentences:
            continue  # Skip empty topics

        full_text = '. '.join(sentences)  # Combine all sentences for the topic
        words = full_text.split()
        chunks = []  # Stores 5-sentence chunks for the topic

        sentence_count = 0
        temp_chunk = []

        # Create chunks with up to 5 sentences
        for word in words:
            temp_chunk.append(word)
            if end_sentence_marker in word:
                sentence_count += 1
            if sentence_count == max_sentences:
                chunks.append(' '.join(temp_chunk))
                temp_chunk = []
                sentence_count = 0

        # Add remaining words in the last chunk
        if temp_chunk:
            chunks.append(' '.join(temp_chunk))

        # Create slides for each chunk
        for chunk_idx, chunk in enumerate(chunks):
            if not chunk.strip():  # Skip blank chunks
                continue

            slide = presentation.slides.add_slide(slide_layout)
            title = slide.shapes.title
            content = slide.placeholders[1]

            # Title for the first slide of the topic
            if first_slide_flag:
                title.text = f"Topic: {topic_keywords[i]}"
                first_slide_flag = False
            else:
                title.text = f"Topic: {topic_keywords[i]} (continued)"

            # Set slide content
            content.text = chunk

            # Adjust font size for the content
            for paragraph in content.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(20)  # Set font size to 20 points

        # Reset the flag for the next topic
        first_slide_flag = True

    # Save the presentation
    presentation.save(output_file)
    print(f"Presentation saved as {output_file}")



# Main Code
if __name__ == "__main__":
    pdf_name = input("Enter the name of the PDF file (with extension): ").strip()
    summary_ratio = input("Enter the summary ratio (e.g., 0.2 for 20%): ").strip()
    n_topics = input("Enter the number of topics to divide into: ").strip()
    
    try:
        summary_ratio = float(summary_ratio)
        n_topics = int(n_topics)
        if not (0 < summary_ratio <= 1):
            raise ValueError("Ratio must be between 0 and 1.")
        if n_topics < 1:
            raise ValueError("Number of topics must be at least 1.")
    except ValueError as e:
        print(f"Invalid input: {e}")
        exit()

    try:
        reader = PdfReader(pdf_name)
    except FileNotFoundError:
        print("The specified PDF file was not found. Please check the file name and path.")
        exit()
    except Exception as e:
        print(f"An error occurred while reading the PDF: {e}")
        exit()

    print(f"The PDF has {len(reader.pages)} pages.\n")

    # Extract text from the entire PDF
    full_text = ""
    for i, page in enumerate(reader.pages):
        try:
            full_text += page.extract_text()
        except Exception as e:
            print(f"An error occurred while extracting text from page {i + 1}: {e}")

    if not full_text.strip():
        print("No text could be extracted from the PDF.")
        exit()

    # Summarize the extracted text
    print("Summarizing the text...")
    summary = summarize_text(full_text, summary_ratio)
    if not summary:
        print("Failed to generate a summary.")
        exit()

    # Divide summary into topics
    print("Dividing summary into topics...")
    topics, topic_keywords = divide_into_topics(summary, n_topics)

    # Create PowerPoint presentation
    create_presentation_with_topics(topics, topic_keywords)
